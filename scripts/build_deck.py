"""Build the SignalLoop case-study deck (LinkedIn carousel) with python-pptx.

Run:  python scripts/build_deck.py   ->  docs/SignalLoop_Deck.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- palette (matches the dashboard) ----------
INK = RGBColor(0x24, 0x1B, 0x4B)
PURPLE = RGBColor(0x6C, 0x5C, 0xE7)
PURPLE_LT = RGBColor(0xA2, 0x9B, 0xFE)
GREEN = RGBColor(0x00, 0xB8, 0x94)
BLUE = RGBColor(0x09, 0x84, 0xE3)
ORANGE = RGBColor(0xE1, 0x70, 0x55)
AMBER = RGBColor(0xFD, 0xCB, 0x6E)
BGLT = RGBColor(0xF5, 0xF2, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x67, 0x90)
INK_SOFT = RGBColor(0xD8, 0xD3, 0xF2)

HEAD = "Segoe UI Semibold"
BODY = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(s, runs, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=BODY, anchor=MSO_ANCHOR.TOP, italic=False, space=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.05))
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if isinstance(runs, str):
        runs = [(runs, {})]
    for i, (line, o) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = o.get("align", align)
        p.space_after = Pt(o.get("space", space))
        r = p.add_run()
        r.text = line
        f = r.font
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold)
        f.italic = o.get("italic", italic)
        f.name = o.get("font", font)
        f.color.rgb = o.get("color", color)
    return tb


def rect(s, x, y, w, h, fill, line=None, line_w=1.0, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def dot(s, x, y, d, color):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def kicker(s, label, color=PURPLE):
    dot(s, 0.85, 0.78, 0.16, color)
    text(s, label, 1.12, 0.66, 6, 0.4, size=13, color=color, bold=True, font=HEAD)


def card(s, x, y, w, h, fill=WHITE, line=RGBColor(0xE7, 0xE2, 0xFB)):
    rect(s, x, y, w, h, fill, line=line, line_w=1.0)


# ============================ SLIDE 1 — TITLE ============================
s = slide(INK)
for i, c in enumerate([PURPLE, GREEN, BLUE, AMBER]):
    dot(s, 0.9 + i * 0.32, 1.4, 0.18, c)
text(s, "SignalLoop", 0.85, 2.25, 11.6, 1.3, size=66, color=WHITE, bold=True, font=HEAD)
text(s, "Turning raw user feedback into an evidence-backed product roadmap.",
     0.9, 3.65, 11, 0.9, size=24, color=PURPLE_LT, font=BODY)
text(s, [("Anuj Kumar", {"bold": True, "color": WHITE}),
         ("  ·  Aspiring AI Product Manager", {"color": INK_SOFT})],
     0.9, 4.7, 11, 0.5, size=17)
text(s, "signalloop.streamlit.app    ·    github.com/sainianuj07/signalloop",
     0.9, 6.7, 11.5, 0.4, size=14, color=INK_SOFT)

# ============================ SLIDE 2 — PROBLEM ============================
s = slide(BGLT)
kicker(s, "THE PROBLEM")
text(s, "Teams drown in feedback — and decide by gut.", 0.85, 1.25, 11.6, 1.0,
     size=34, color=INK, bold=True, font=HEAD)
text(s, [("Early-stage SaaS teams collect feedback everywhere — app reviews, tickets, NPS — "
          "but nobody has time to read it all.", {"space": 10}),
         ("So roadmap calls default to the loudest customer or a hunch. The signal is in there; "
          "it just never gets found.", {})],
     0.9, 2.5, 7.0, 2.4, size=18, color=GREY)
card(s, 8.4, 2.45, 4.1, 2.6, fill=WHITE)
text(s, "216", 8.4, 2.7, 4.1, 1.0, size=64, color=PURPLE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
text(s, "real reviews", 8.4, 3.8, 4.1, 0.4, size=16, color=GREY, align=PP_ALIGN.CENTER)
text(s, "0 PMs with time to read them", 8.4, 4.25, 4.1, 0.5, size=15, color=ORANGE,
     bold=True, align=PP_ALIGN.CENTER)

# ============================ SLIDE 3 — THE IDEA (pipeline) ============================
s = slide(WHITE)
kicker(s, "THE IDEA")
text(s, "Close the loop — feedback to decision.", 0.85, 1.25, 11.6, 0.9,
     size=34, color=INK, bold=True, font=HEAD)
stages = [("Ingest", "Reviews & tickets\ninto one place", PURPLE),
          ("Understand", "LLM labels every\nitem (type, severity)", BLUE),
          ("Prioritize", "Cluster into themes\n→ rank with RICE", GREEN),
          ("Act", "PRDs, decisions,\nweekly briefs", ORANGE)]
cw, gap, y0 = 2.72, 0.3, 3.0
for i, (t, d, c) in enumerate(stages):
    x = 0.85 + i * (cw + gap)
    card(s, x, y0, cw, 2.3)
    rect(s, x, y0, cw, 0.12, c, radius=False)
    text(s, t, x + 0.1, y0 + 0.35, cw - 0.2, 0.5, size=20, color=INK, bold=True, font=HEAD)
    text(s, d, x + 0.1, y0 + 1.0, cw - 0.2, 1.1, size=14, color=GREY)
    if i < 3:
        text(s, "→", x + cw - 0.02, y0 + 0.85, 0.34, 0.5, size=24, color=PURPLE_LT,
             bold=True, align=PP_ALIGN.CENTER)
text(s, "Built end-to-end on 216 live Notion reviews — not a toy dataset.",
     0.9, 5.9, 11, 0.5, size=15, color=GREY, italic=True)

# ============================ SLIDE 4 — UNDERSTAND ============================
s = slide(BGLT)
kicker(s, "UNDERSTAND", BLUE)
text(s, "Every review becomes structured data.", 0.85, 1.25, 11.6, 0.9,
     size=34, color=INK, bold=True, font=HEAD)
text(s, "An LLM turns messy text into clean JSON the product can actually use — labelled by "
        "type, product area, sentiment, and severity. Batched, retried, with an offline mock mode.",
     0.9, 2.45, 11.4, 1.0, size=18, color=GREY)
chips = ["bug", "feature_request", "ux_friction", "pricing", "churn_risk", "question", "praise"]
cx = 0.9
for ch in chips:
    cwid = 0.32 + len(ch) * 0.115
    rect(s, cx, 3.75, cwid, 0.55, PURPLE_LT)
    text(s, ch, cx, 3.84, cwid, 0.4, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
    cx += cwid + 0.2
for i, (n, l, c) in enumerate([("216", "reviews labelled", PURPLE),
                               ("8", "feedback types", BLUE),
                               ("1-4", "severity scale", GREEN)]):
    x = 0.9 + i * 4.05
    card(s, x, 4.8, 3.7, 1.5)
    text(s, n, x, 4.95, 3.7, 0.8, size=42, color=c, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
    text(s, l, x, 5.8, 3.7, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)

# ============================ SLIDE 5 — EVAL (HERO) ============================
s = slide(INK)
dot(s, 0.85, 0.78, 0.16, GREEN)
text(s, "THE EVAL SPINE", 1.12, 0.66, 6, 0.4, size=13, color=GREEN, bold=True, font=HEAD)
text(s, "“How do I know the AI is right?”", 0.85, 1.25, 11.6, 0.9,
     size=34, color=WHITE, bold=True, font=HEAD)
text(s, [("I hand-labelled a golden set, then built an ", {}),
         ("independent LLM judge", {"bold": True, "color": PURPLE_LT}),
         (" — OpenAI’s model grading Meta’s, so it can’t just agree with itself. "
          "I only trusted it after it matched my human labels.", {})],
     0.9, 2.4, 7.1, 2.2, size=18, color=INK_SOFT)
card(s, 8.3, 2.35, 4.2, 1.5, fill=RGBColor(0x32, 0x2A, 0x5E), line=None)
text(s, "~96%", 8.3, 2.5, 4.2, 0.8, size=46, color=GREEN, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
text(s, "judge-estimated type accuracy", 8.3, 3.35, 4.2, 0.4, size=13, color=INK_SOFT,
     align=PP_ALIGN.CENTER)
rect(s, 0.9, 4.95, 11.55, 1.75, RGBColor(0x32, 0x2A, 0x5E))
text(s, "It caught a bias I’d never have spotted by eye:", 1.2, 5.15, 11, 0.45,
     size=16, color=AMBER, bold=True)
text(s, [("The classifier tagged ", {}), ("angry", {"italic": True, "color": WHITE}),
         (" reviews as ", {}), ("churn_risk", {"bold": True, "color": PURPLE_LT}),
         (" even when nobody said they’d leave. I tightened the prompt (anger ≠ churn), "
          "re-ran the eval, and the bad labels dropped 14 → 9. Measure → fix → measure.", {})],
     1.2, 5.6, 11.1, 1.0, size=15, color=INK_SOFT)

# ============================ SLIDE 6 — PRIORITIZE ============================
s = slide(WHITE)
kicker(s, "PRIORITIZE", GREEN)
text(s, "From themes to a ranked roadmap.", 0.85, 1.25, 11.6, 0.9,
     size=34, color=INK, bold=True, font=HEAD)
text(s, "Reviews are embedded, clustered into named themes, then scored with RICE = "
        "(Reach × Impact × Confidence) / Effort — so priority is grounded in data, not volume.",
     0.9, 2.45, 11.4, 1.0, size=18, color=GREY)
rows = [("App Performance & Stability", 34.8, 1.0),
        ("UI/UX & Mobile Hurdles", 21.2, 0.61),
        ("Functionality & Usability", 19.4, 0.56),
        ("Intrusive AI Features", 14.0, 0.40),
        ("False Offline Status", 12.8, 0.37)]
y = 3.7
for name, score, frac in rows:
    text(s, name, 0.9, y, 4.6, 0.4, size=15, color=INK, bold=True)
    rect(s, 5.6, y + 0.05, 6.0 * frac, 0.32, PURPLE)
    rect(s, 5.6 + 6.0 * frac + 0.06, y + 0.02, 0.9, 0.36, WHITE)
    text(s, str(score), 5.6 + 6.0 * frac + 0.1, y + 0.03, 1.0, 0.34, size=14, color=PURPLE,
         bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.62

# ============================ SLIDE 7 — PRD ============================
s = slide(BGLT)
kicker(s, "ACT · I", ORANGE)
text(s, "PRDs that can’t hallucinate quotes.", 0.85, 1.25, 11.6, 0.9,
     size=34, color=INK, bold=True, font=HEAD)
text(s, [("The generator drafts a full PRD where ", {}),
         ("every claim cites a real review ID", {"bold": True, "color": INK}),
         (" like [#149]. A checker verifies each citation exists, and the evidence quotes are "
          "pulled straight from the database — not the model.", {})],
     0.9, 2.45, 7.0, 2.2, size=18, color=GREY)
text(s, "So the model picks which reviews to cite; my code supplies the actual words. "
        "A faked quote is impossible by construction.",
     0.9, 4.5, 7.0, 1.2, size=16, color=INK, italic=True)
card(s, 8.3, 2.5, 4.2, 2.8, WHITE)
text(s, "Citation check", 8.3, 2.7, 4.2, 0.4, size=15, color=GREY, bold=True, align=PP_ALIGN.CENTER)
text(s, "15 / 15", 8.3, 3.15, 4.2, 0.9, size=52, color=GREEN, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
text(s, "citations valid", 8.3, 4.15, 4.2, 0.4, size=15, color=GREY, align=PP_ALIGN.CENTER)
text(s, "0 hallucinated", 8.3, 4.6, 4.2, 0.45, size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# ============================ SLIDE 8 — BOARDROOM ============================
s = slide(WHITE)
kicker(s, "ACT · II", ORANGE)
text(s, "Prioritization is a negotiation, not a formula.", 0.85, 1.25, 11.6, 0.9,
     size=32, color=INK, bold=True, font=HEAD)
text(s, "Three agents debate the roadmap from their own incentives; a Head-of-Product agent "
        "makes the call — and records who got overruled.",
     0.9, 2.4, 11.4, 0.9, size=18, color=GREY)
agents = [("Growth PM", "reach · retention", PURPLE),
          ("Engineering", "effort · risk", BLUE),
          ("Finance", "cost · ROI", GREEN)]
for i, (t, d, c) in enumerate(agents):
    x = 0.9 + i * 3.0
    card(s, x, 3.5, 2.75, 1.5)
    rect(s, x, 3.5, 0.12, 1.5, c, radius=False)
    text(s, t, x + 0.25, 3.7, 2.4, 0.4, size=18, color=INK, bold=True, font=HEAD)
    text(s, d, x + 0.25, 4.2, 2.4, 0.5, size=14, color=GREY)
text(s, "→", 9.7, 3.95, 0.6, 0.6, size=30, color=PURPLE_LT, bold=True, align=PP_ALIGN.CENTER)
rect(s, 10.3, 3.5, 2.2, 1.5, INK)
text(s, "Head of Product", 10.3, 3.72, 2.2, 0.4, size=15, color=WHITE, bold=True,
     font=HEAD, align=PP_ALIGN.CENTER)
text(s, "decides + logs dissent", 10.3, 4.2, 2.2, 0.5, size=12, color=PURPLE_LT, align=PP_ALIGN.CENTER)
text(s, "“I’m knowingly setting aside Engineering’s concern on False Offline Status — "
        "and I’m comfortable carrying that risk.”",
     0.9, 5.45, 11.5, 0.9, size=15, color=GREY, italic=True)

# ============================ SLIDE 9 — DRIFT ============================
s = slide(BGLT)
kicker(s, "ACT · III", ORANGE)
text(s, "Watching what’s changing, not just what is.", 0.85, 1.25, 11.6, 0.9,
     size=34, color=INK, bold=True, font=HEAD)
text(s, "A static dashboard shows the state of feedback. Drift compares this week to last and "
        "flags the spikes — the emerging issues a ranked table buries.",
     0.9, 2.45, 7.0, 1.6, size=18, color=GREY)
card(s, 8.2, 2.5, 4.3, 2.9, WHITE)
text(s, "False Offline Status", 8.2, 2.75, 4.3, 0.4, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, [("1 ", {"size": 40, "color": GREY, "bold": True}),
         (" → ", {"size": 30, "color": PURPLE_LT, "bold": True}),
         (" 6", {"size": 56, "color": ORANGE, "bold": True})],
     8.2, 3.3, 4.3, 1.1, font=HEAD, align=PP_ALIGN.CENTER)
text(s, "complaints, week over week", 8.2, 4.5, 4.3, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)
text(s, "+ an auto-written “State of Feedback” brief", 8.2, 4.9, 4.3, 0.4, size=13,
     color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
text(s, "6× in a week often means a fresh regression — exactly what you want flagged early.",
     0.9, 4.3, 7.0, 0.9, size=16, color=INK, italic=True)

# ============================ SLIDE 10 — LEARNED ============================
s = slide(WHITE)
kicker(s, "WHAT I LEARNED")
text(s, "The honest part.", 0.85, 1.25, 11.6, 0.9, size=34, color=INK, bold=True, font=HEAD)
lessons = [("Rate limits were the real boss fight",
            "Not model quality. I built provider fallback, cooldowns, and an embedding cache to survive free tiers.", ORANGE),
           ("Not everything clusters cleanly",
            "A silhouette score (~0.04) showed the feedback is a continuum — so “how many themes” is a product judgment, not a metric.", BLUE),
           ("Models are confidently wrong",
            "Every label came back 80–100% sure, including the wrong ones — which is exactly why an independent eval matters.", GREEN)]
y = 2.55
for t, d, c in lessons:
    rect(s, 0.9, y, 0.14, 1.25, c, radius=False)
    text(s, t, 1.25, y, 11.0, 0.5, size=20, color=INK, bold=True, font=HEAD)
    text(s, d, 1.25, y + 0.5, 11.0, 0.8, size=15, color=GREY)
    y += 1.45

# ============================ SLIDE 11 — CLOSE ============================
s = slide(INK)
text(s, "Built end-to-end. Live & open-source.", 0.9, 2.4, 11.5, 1.0,
     size=38, color=WHITE, bold=True, font=HEAD)
text(s, "Ingestion → classification → evaluation → themes → RICE → PRDs → boardroom → drift.",
     0.9, 3.5, 11.5, 0.6, size=18, color=PURPLE_LT)
for i, (label, val) in enumerate([("Live demo", "signalloop.streamlit.app"),
                                   ("Code", "github.com/sainianuj07/signalloop"),
                                   ("Contact", "anuj.product.iitr@gmail.com")]):
    y = 4.5 + i * 0.62
    text(s, label, 0.9, y, 2.0, 0.4, size=15, color=GREY, bold=True)
    text(s, val, 3.0, y, 9.0, 0.4, size=16, color=WHITE)
text(s, "Open to AI Product Management roles.", 0.9, 6.55, 11.5, 0.5, size=18,
     color=GREEN, bold=True, font=HEAD)

# ---------- save ----------
out = Path("docs/SignalLoop_Deck.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"Saved {len(prs.slides._sldIdLst)} slides -> {out}")
